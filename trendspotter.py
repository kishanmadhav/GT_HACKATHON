"""
TrendSpotter - Automated AdTech Insight Engine

Drops any marketing data file and generates professional reports with:
- Automatic column detection (works with messy exports)
- Anomaly detection using Isolation Forest
- AI-powered insights via GPT-4o
- Charts and visualizations

Usage: python trendspotter.py [INPUT_FILE] [--pdf | --pptx]
"""

import os
import sys
import time
from datetime import datetime
from pathlib import Path
import polars as pl
import numpy as np
from sklearn.ensemble import IsolationForest
from openai import OpenAI
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import re

API_KEY = os.getenv("OPENAI_API_KEY", "")
INPUT_FILE = "data/sample/ad_performance.csv"
OUTPUT_DIR = "output"
CHARTS_DIR = "output/charts"

BLUE = '#2E86AB'
PURPLE = '#A23B72'
GREEN = '#28A745'
RED = '#DC3545'
ORANGE = '#FFC107'

class DataNormalizer:
    """Figures out what each column means, even if the headers are weird"""
    
    # These patterns catch most common naming conventions in marketing exports
    COLUMN_PATTERNS = {
        'date': [r'date', r'day', r'time', r'timestamp', r'created', r'period', r'dt', r'report'],
        'campaign_id': [r'campaign_id', r'camp_id', r'adgroup_id', r'ad_id', r'^id$', r'cid'],
        'campaign_name': [r'campaign_name', r'name', r'title', r'ad_name', r'adgroup', r'ad_group', r'campaign$'],
        'impressions': [r'impression', r'impr', r'views', r'imp', r'shows', r'reach', r'display'],
        'clicks': [r'click', r'clk', r'visits', r'sessions', r'hits', r'tap'],
        'conversions': [r'conversion', r'conv', r'purchase', r'order', r'lead', r'action', r'signup', r'install'],
        'spend': [r'spend', r'cost', r'expense', r'budget', r'amount_spent', r'media_cost', r'ad_cost'],
        'revenue': [r'revenue', r'income', r'sales', r'value', r'earning', r'gmv', r'amount', r'total'],
        'region': [r'region', r'country', r'location', r'geo', r'territory', r'market', r'area', r'state', r'city'],
        'platform': [r'platform', r'source', r'channel', r'network', r'medium', r'publisher', r'traffic']
    }
    
    def __init__(self):
        self.mapping = {}
        self.metadata = {}
    
    def _match_column(self, col_name):
        col_lower = col_name.lower().strip()
        for standard, patterns in self.COLUMN_PATTERNS.items():
            for p in patterns:
                if re.search(p, col_lower):
                    return standard
        return None
    
    def _infer_type(self, df, col):
        dtype = df[col].dtype
        if dtype == pl.Date or dtype == pl.Datetime:
            return 'date'
        if dtype == pl.Utf8:
            sample = str(df[col].drop_nulls().head(1)[0]) if len(df[col].drop_nulls()) > 0 else ''
            if re.match(r'\d{4}-\d{2}-\d{2}', sample) or re.match(r'\d{2}/\d{2}/\d{4}', sample):
                return 'date'
            return 'text'
        if dtype in [pl.Int64, pl.Int32, pl.Float64, pl.Float32]:
            mean = df[col].mean()
            if mean and mean > 100:
                return 'metric_large'
            return 'metric_small'
        return 'unknown'
    
    def normalize(self, df):
        """Normalize any dataframe to standard schema"""
        self.mapping = {}
        mapped = set()
        unmapped = []
        
        # Pattern matching
        for col in df.columns:
            std = self._match_column(col)
            if std and std not in mapped:
                self.mapping[col] = std
                mapped.add(std)
            else:
                unmapped.append(col)
        
        # Type inference for remaining
        for col in unmapped[:]:
            col_type = self._infer_type(df, col)
            if col_type == 'date' and 'date' not in mapped:
                self.mapping[col] = 'date'
                mapped.add('date')
                unmapped.remove(col)
            elif col_type == 'metric_large' and 'impressions' not in mapped:
                self.mapping[col] = 'impressions'
                mapped.add('impressions')
                unmapped.remove(col)
        
        # Rename mapped columns
        result = df.rename({k: v for k, v in self.mapping.items()})
        cols = set(result.columns)
        synthetic = []
        
        # Create missing columns
        if 'date' not in cols:
            result = result.with_columns(pl.lit(datetime.now().strftime('%Y-%m-%d')).alias('date'))
            synthetic.append('date')
        elif result['date'].dtype == pl.Utf8:
            try:
                result = result.with_columns(pl.col('date').str.to_date().alias('date'))
            except:
                pass
        
        if 'campaign_name' not in cols:
            if 'campaign_id' in cols:
                result = result.with_columns(pl.col('campaign_id').cast(pl.Utf8).alias('campaign_name'))
            else:
                result = result.with_columns(pl.lit('Campaign_1').alias('campaign_name'))
            synthetic.append('campaign_name')
        
        if 'campaign_id' not in cols:
            result = result.with_columns(pl.arange(1, len(result)+1).cast(pl.Utf8).alias('campaign_id'))
            synthetic.append('campaign_id')
        
        # Numeric defaults
        for field, default in [('impressions', 0), ('clicks', 0), ('conversions', 0), 
                               ('spend', 0.0), ('revenue', 0.0)]:
            if field not in result.columns:
                # Try to use an unmapped numeric column
                found = False
                for col in unmapped[:]:
                    if col in result.columns and result[col].dtype in [pl.Int64, pl.Int32, pl.Float64, pl.Float32]:
                        result = result.rename({col: field})
                        unmapped.remove(col)
                        found = True
                        break
                if not found:
                    result = result.with_columns(pl.lit(default).alias(field))
                    synthetic.append(field)
        
        if 'region' not in result.columns:
            result = result.with_columns(pl.lit('Unknown').alias('region'))
            synthetic.append('region')
        if 'platform' not in result.columns:
            result = result.with_columns(pl.lit('Unknown').alias('platform'))
            synthetic.append('platform')
        
        self.metadata = {
            'mapping': self.mapping,
            'synthetic': synthetic,
            'original_cols': df.columns
        }
        
        return result


def load_data(file_path):
    """Load any supported file format and normalize the columns"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Can't find the file: {file_path}")
    
    ext = file_path.lower().split('.')[-1]
    
    try:
        if ext == 'csv':
            df = pl.read_csv(file_path, try_parse_dates=True)
        elif ext == 'json':
            try:
                df = pl.read_json(file_path)
            except:
                df = pl.read_ndjson(file_path)
        elif ext == 'ndjson':
            df = pl.read_ndjson(file_path)
        elif ext == 'parquet':
            df = pl.read_parquet(file_path)
        elif ext in ['sqlite', 'db', 'sqlite3']:
            import sqlite3
            conn = sqlite3.connect(file_path)
            cursor = conn.cursor()
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table' LIMIT 1")
            result = cursor.fetchone()
            if not result:
                conn.close()
                raise ValueError(f"No tables found in database: {file_path}")
            table_name = result[0]
            df = pl.read_database(f"SELECT * FROM {table_name}", conn)
            conn.close()
        else:
            raise ValueError(f"Unsupported file type: {ext}. Use CSV, JSON, NDJSON, Parquet, or SQLite.")
    except Exception as e:
        raise RuntimeError(f"Failed to read {file_path}: {str(e)}")
    
    if df.is_empty():
        raise ValueError(f"File is empty: {file_path}")
    
    normalizer = DataNormalizer()
    normalized = normalizer.normalize(df)
    
    return normalized, normalizer.metadata


def generate_charts(df, dates, campaigns, metadata):
    """Build all the charts we need for the report"""
    os.makedirs(CHARTS_DIR, exist_ok=True)
    
    # Daily aggregation
    daily = df.group_by('date').agg([
        pl.col('impressions').sum(),
        pl.col('clicks').sum(),
        pl.col('revenue').sum(),
        pl.col('spend').sum()
    ]).sort('date')
    
    d_dates = [str(d) for d in daily['date'].to_list()]  # Convert to strings
    d_imp = daily['impressions'].to_list()
    d_rev = daily['revenue'].to_list()
    d_spend = daily['spend'].to_list()
    d_clicks = daily['clicks'].to_list()
    
    # Format date labels
    def format_date(d):
        if len(d) >= 10:
            return d[5:10]  # MM-DD
        return d[-5:] if len(d) >= 5 else d
    
    # Chart 1: Dashboard (2x2)
    fig, axes = plt.subplots(2, 2, figsize=(12, 8))
    fig.suptitle('Performance Dashboard', fontsize=16, fontweight='bold', color=BLUE)
    
    step = max(1, len(d_dates) // 5)  # Show ~5 x-labels
    
    # Impressions trend
    axes[0, 0].fill_between(range(len(d_dates)), d_imp, alpha=0.3, color=BLUE)
    axes[0, 0].plot(d_imp, color=BLUE, linewidth=2, marker='o', markersize=4)
    axes[0, 0].set_title('Daily Impressions', fontweight='bold')
    axes[0, 0].set_xticks(range(0, len(d_dates), step))
    axes[0, 0].set_xticklabels([format_date(d_dates[i]) for i in range(0, len(d_dates), step)], rotation=45)
    axes[0, 0].grid(True, alpha=0.3)
    
    # Revenue vs Spend
    x = range(len(d_dates))
    width = 0.35
    axes[0, 1].bar([i - width/2 for i in x], d_rev, width, label='Revenue', color=GREEN)
    axes[0, 1].bar([i + width/2 for i in x], d_spend, width, label='Spend', color=RED)
    axes[0, 1].set_title('Revenue vs Spend', fontweight='bold')
    axes[0, 1].legend()
    axes[0, 1].set_xticks(range(0, len(d_dates), step))
    axes[0, 1].set_xticklabels([format_date(d_dates[i]) for i in range(0, len(d_dates), step)], rotation=45)
    axes[0, 1].grid(True, alpha=0.3)
    
    # Clicks trend
    axes[1, 0].fill_between(range(len(d_dates)), d_clicks, alpha=0.3, color=PURPLE)
    axes[1, 0].plot(d_clicks, color=PURPLE, linewidth=2, marker='o', markersize=4)
    axes[1, 0].set_title('Daily Clicks', fontweight='bold')
    axes[1, 0].set_xticks(range(0, len(d_dates), step))
    axes[1, 0].set_xticklabels([format_date(d_dates[i]) for i in range(0, len(d_dates), step)], rotation=45)
    axes[1, 0].grid(True, alpha=0.3)
    
    # ROI trend
    d_roi = [((r-s)/s*100) if s > 0 else 0 for r, s in zip(d_rev, d_spend)]
    axes[1, 1].plot(d_roi, color=ORANGE, linewidth=2, marker='s', markersize=4)
    axes[1, 1].axhline(y=np.mean(d_roi), color='gray', linestyle='--', label=f'Avg: {np.mean(d_roi):.0f}%')
    axes[1, 1].set_title('Daily ROI %', fontweight='bold')
    axes[1, 1].legend()
    axes[1, 1].set_xticks(range(0, len(d_dates), step))
    axes[1, 1].set_xticklabels([format_date(d_dates[i]) for i in range(0, len(d_dates), step)], rotation=45)
    axes[1, 1].grid(True, alpha=0.3)
    
    plt.tight_layout()
    plt.savefig(f'{CHARTS_DIR}/dashboard.png', dpi=150, bbox_inches='tight')
    plt.close()
    
    # Chart 2: Campaign Revenue Bar Chart
    fig, ax = plt.subplots(figsize=(10, 6))
    camp_names = campaigns['campaign_name'].to_list()
    camp_rev = campaigns['revenue'].to_list()
    colors = [BLUE, PURPLE, GREEN, ORANGE, RED] * (len(camp_names) // 5 + 1)
    bars = ax.bar(camp_names[:10], camp_rev[:10], color=colors[:min(10, len(camp_names))])
    ax.set_title('Revenue by Campaign', fontsize=14, fontweight='bold', color=BLUE)
    ax.set_ylabel('Revenue ($)')
    for bar, val in zip(bars, camp_rev[:10]):
        if val > 0:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(camp_rev)*0.02, 
                    f'${val:,.0f}', ha='center', va='bottom', fontsize=9)
    plt.xticks(rotation=30, ha='right')
    plt.grid(True, alpha=0.3, axis='y')
    plt.subplots_adjust(bottom=0.18)
    plt.savefig(f'{CHARTS_DIR}/campaign_revenue.png', dpi=150, bbox_inches='tight')
    plt.close()
    
    # Chart 3: Platform Pie Chart
    if 'platform' in df.columns:
        platforms = df.group_by('platform').agg(pl.col('revenue').sum()).sort('revenue', descending=True)
        fig, ax = plt.subplots(figsize=(8, 6))
        ax.pie(platforms['revenue'].to_list(), labels=platforms['platform'].to_list(),
               autopct='%1.1f%%', colors=[BLUE, PURPLE, GREEN, ORANGE, RED], startangle=90)
        ax.set_title('Revenue by Platform', fontsize=14, fontweight='bold', color=BLUE)
        plt.tight_layout()
        plt.savefig(f'{CHARTS_DIR}/platform_pie.png', dpi=150, bbox_inches='tight')
        plt.close()
    
    return {
        'dashboard': f'{CHARTS_DIR}/dashboard.png',
        'campaign': f'{CHARTS_DIR}/campaign_revenue.png',
        'platform': f'{CHARTS_DIR}/platform_pie.png'
    }


def generate_pdf(data, input_filename):
    """Build an executive-ready PDF report with beautiful formatting"""
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    report_title = Path(input_filename).stem.replace('_', ' ').replace('-', ' ').title()
    kpis = data['kpis']
    
    # Calculate key insights for highlighting
    top_campaign = data['campaigns'].row(0, named=True)
    top_roi = ((top_campaign['revenue'] - top_campaign['spend']) / top_campaign['spend'] * 100) if top_campaign['spend'] > 0 else 0
    
    # Page 1: Cover Page
    pdf.add_page()
    
    # Gradient-style header with darker overlay
    pdf.set_fill_color(30, 60, 90)
    pdf.rect(0, 0, 210, 297, 'F')
    
    # Accent stripe
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(0, 100, 210, 8, 'F')
    
    # Main title
    pdf.set_font('Helvetica', 'B', 36)
    pdf.set_text_color(255, 255, 255)
    pdf.set_xy(20, 50)
    pdf.cell(0, 15, report_title)
    
    pdf.set_font('Helvetica', '', 18)
    pdf.set_xy(20, 70)
    pdf.cell(0, 10, 'Performance Analytics Report')
    
    # Date range badge
    pdf.set_xy(20, 120)
    pdf.set_font('Helvetica', 'B', 14)
    pdf.set_text_color(200, 220, 240)
    pdf.cell(0, 8, f"{data['dates'][0]}  to  {data['dates'][-1]}")
    
    # Hero metrics on cover
    pdf.set_y(150)
    hero_metrics = [
        ('Total Revenue', f"${kpis['revenue']:,.0f}", GREEN),
        ('Return on Investment', f"{kpis['roi']:.1f}%", BLUE if kpis['roi'] > 0 else RED),
        ('Anomalies Detected', f"{len(data['anomalies'])}", ORANGE if data['anomalies'] else GREEN),
    ]
    
    for i, (label, value, color) in enumerate(hero_metrics):
        y_pos = 150 + i * 35
        pdf.set_fill_color(40, 70, 100)
        pdf.rect(20, y_pos, 170, 28, 'F')
        
        # Color accent bar
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        pdf.set_fill_color(r, g, b)
        pdf.rect(20, y_pos, 4, 28, 'F')
        
        pdf.set_xy(30, y_pos + 4)
        pdf.set_font('Helvetica', '', 10)
        pdf.set_text_color(180, 200, 220)
        pdf.cell(0, 5, label.upper())
        
        pdf.set_xy(30, y_pos + 12)
        pdf.set_font('Helvetica', 'B', 18)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 8, value)
    
    # Footer
    pdf.set_xy(20, 265)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(120, 140, 160)
    pdf.cell(0, 5, f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}")
    pdf.set_xy(20, 272)
    pdf.cell(0, 5, "Powered by TrendSpotter AI Analytics")
    
    # Page 2: Key Metrics Dashboard
    pdf.add_page()
    pdf.set_fill_color(255, 255, 255)
    pdf.rect(0, 0, 210, 297, 'F')
    
    # Section header with accent
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(0, 0, 210, 3, 'F')
    
    pdf.set_xy(15, 12)
    pdf.set_font('Helvetica', 'B', 20)
    pdf.set_text_color(30, 60, 90)
    pdf.cell(0, 10, 'Key Performance Indicators')
    
    # KPI Cards in 2x3 grid
    kpi_cards = [
        ('IMPRESSIONS', f"{kpis['impressions']:,}", 'Total ad views', BLUE),
        ('CLICKS', f"{kpis['clicks']:,}", f"CTR: {kpis['ctr']:.2f}%", PURPLE),
        ('CONVERSIONS', f"{kpis['conversions']:,}", 'Completed actions', GREEN),
        ('REVENUE', f"${kpis['revenue']:,.0f}", 'Total earnings', GREEN),
        ('AD SPEND', f"${kpis['spend']:,.0f}", 'Investment', ORANGE),
        ('ROI', f"{kpis['roi']:.1f}%", 'Return on investment', GREEN if kpis['roi'] > 50 else ORANGE if kpis['roi'] > 0 else RED),
    ]
    
    for i, (label, value, subtitle, color) in enumerate(kpi_cards):
        col, row = i % 3, i // 3
        x = 15 + col * 62
        y = 32 + row * 45
        
        # Card background
        pdf.set_fill_color(248, 250, 252)
        pdf.rect(x, y, 58, 40, 'F')
        
        # Top accent bar
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        pdf.set_fill_color(r, g, b)
        pdf.rect(x, y, 58, 3, 'F')
        
        # Label
        pdf.set_xy(x + 4, y + 7)
        pdf.set_font('Helvetica', 'B', 8)
        pdf.set_text_color(100, 110, 120)
        pdf.cell(50, 4, label)
        
        # Value
        pdf.set_xy(x + 4, y + 15)
        pdf.set_font('Helvetica', 'B', 16)
        pdf.set_text_color(30, 40, 50)
        pdf.cell(50, 8, value)
        
        # Subtitle
        pdf.set_xy(x + 4, y + 28)
        pdf.set_font('Helvetica', '', 7)
        pdf.set_text_color(130, 140, 150)
        pdf.cell(50, 4, subtitle)
    
    # Key Insight Highlight Box
    pdf.set_y(128)
    pdf.set_fill_color(240, 249, 255)
    pdf.rect(15, 128, 180, 35, 'F')
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(15, 128, 4, 35, 'F')
    
    pdf.set_xy(25, 132)
    pdf.set_font('Helvetica', 'B', 10)
    pdf.set_text_color(46, 134, 171)
    pdf.cell(0, 5, 'KEY INSIGHT')
    
    pdf.set_xy(25, 140)
    pdf.set_font('Helvetica', '', 11)
    pdf.set_text_color(40, 50, 60)
    insight_text = f"Top performer: {top_campaign['campaign_name']} generated ${top_campaign['revenue']:,.0f} revenue with {top_roi:.0f}% ROI"
    pdf.multi_cell(165, 6, insight_text)
    
    # Executive Summary Section
    pdf.set_y(172)
    pdf.set_font('Helvetica', 'B', 14)
    pdf.set_text_color(30, 60, 90)
    pdf.cell(0, 8, 'Executive Summary')
    
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(15, 182, 40, 1, 'F')
    
    pdf.set_xy(15, 188)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(50, 60, 70)
    pdf.multi_cell(180, 5.5, data['summary'])
    
    # Anomalies Alert Section
    if data['anomalies']:
        pdf.set_y(pdf.get_y() + 8)
        pdf.set_fill_color(255, 245, 245)
        alert_y = pdf.get_y()
        pdf.rect(15, alert_y, 180, 8 + len(data['anomalies'][:4]) * 10, 'F')
        pdf.set_fill_color(220, 53, 69)
        pdf.rect(15, alert_y, 4, 8 + len(data['anomalies'][:4]) * 10, 'F')
        
        pdf.set_xy(25, alert_y + 3)
        pdf.set_font('Helvetica', 'B', 10)
        pdf.set_text_color(180, 40, 50)
        pdf.cell(0, 5, f'ANOMALIES DETECTED ({len(data["anomalies"])})')
        
        for i, a in enumerate(data['anomalies'][:4]):
            pdf.set_xy(25, alert_y + 12 + i * 10)
            pdf.set_font('Helvetica', '', 9)
            pdf.set_text_color(80, 50, 50)
            pdf.cell(0, 5, f"* {a['desc'][:70]}")
    
    # Page 3: Performance Charts
    pdf.add_page()
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(0, 0, 210, 3, 'F')
    
    pdf.set_xy(15, 12)
    pdf.set_font('Helvetica', 'B', 20)
    pdf.set_text_color(30, 60, 90)
    pdf.cell(0, 10, 'Performance Dashboard')
    
    if os.path.exists(data['charts']['dashboard']):
        pdf.image(data['charts']['dashboard'], x=8, y=28, w=194)
    
    # Page 4: Campaign Breakdown
    pdf.add_page()
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(0, 0, 210, 3, 'F')
    
    pdf.set_xy(15, 12)
    pdf.set_font('Helvetica', 'B', 20)
    pdf.set_text_color(30, 60, 90)
    pdf.cell(0, 10, 'Campaign Performance')
    
    if os.path.exists(data['charts']['campaign']):
        pdf.image(data['charts']['campaign'], x=15, y=28, w=180)
    
    # Campaign Table with better styling - moved down to avoid overlap
    pdf.set_y(118)
    pdf.set_font('Helvetica', 'B', 12)
    pdf.set_text_color(30, 60, 90)
    pdf.cell(0, 8, 'Performance by Campaign')
    pdf.ln(10)
    
    # Table header
    pdf.set_fill_color(30, 60, 90)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font('Helvetica', 'B', 9)
    cols = [('Campaign', 52), ('Impressions', 30), ('Clicks', 26), ('Revenue', 32), ('Spend', 30), ('ROI', 20)]
    x = 15
    for name, w in cols:
        pdf.set_xy(x, pdf.get_y())
        pdf.cell(w, 8, name, 0, 0, 'C', True)
        x += w
    pdf.ln()
    
    pdf.set_font('Helvetica', '', 9)
    for i, row in enumerate(data['campaigns'].head(8).iter_rows(named=True)):
        pdf.set_fill_color(248, 250, 252) if i % 2 == 0 else pdf.set_fill_color(255, 255, 255)
        roi = ((row['revenue'] - row['spend']) / row['spend'] * 100) if row['spend'] > 0 else 0
        
        # Highlight top performer
        if i == 0:
            pdf.set_fill_color(240, 249, 255)
        
        pdf.set_text_color(40, 50, 60)
        vals = [
            row['campaign_name'][:18],
            f"{row['impressions']:,}",
            f"{row['clicks']:,}",
            f"${row['revenue']:,.0f}",
            f"${row['spend']:,.0f}",
            f"{roi:.0f}%"
        ]
        x = 15
        for (_, w), v in zip(cols, vals):
            pdf.set_xy(x, pdf.get_y())
            pdf.cell(w, 7, v, 0, 0, 'C', True)
            x += w
        pdf.ln()
    
    # Page 5: AI Recommendations
    pdf.add_page()
    pdf.set_fill_color(46, 134, 171)
    pdf.rect(0, 0, 210, 3, 'F')
    
    pdf.set_xy(15, 12)
    pdf.set_font('Helvetica', 'B', 20)
    pdf.set_text_color(30, 60, 90)
    pdf.cell(0, 10, 'Strategic Recommendations')
    
    pdf.set_xy(15, 24)
    pdf.set_font('Helvetica', '', 10)
    pdf.set_text_color(100, 110, 120)
    pdf.cell(0, 5, 'AI-powered insights to optimize your campaigns')
    
    # Recommendation cards
    rec_colors = [GREEN, BLUE, PURPLE, ORANGE]
    y_start = 38
    
    for i, rec in enumerate(data['recommendations'][:4]):
        if not rec or len(rec.strip()) < 3:
            continue
        
        y_pos = y_start + i * 32
        color = rec_colors[i % 4]
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        
        # Card background
        pdf.set_fill_color(250, 252, 255)
        pdf.rect(15, y_pos, 180, 26, 'F')
        
        # Left accent
        pdf.set_fill_color(r, g, b)
        pdf.rect(15, y_pos, 4, 26, 'F')
        
        # Number badge
        pdf.set_fill_color(r, g, b)
        pdf.rect(24, y_pos + 4, 18, 18, 'F')
        pdf.set_xy(24, y_pos + 8)
        pdf.set_font('Helvetica', 'B', 12)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(18, 8, str(i + 1), 0, 0, 'C')
        
        # Recommendation text
        pdf.set_xy(48, y_pos + 6)
        pdf.set_font('Helvetica', '', 10)
        pdf.set_text_color(40, 50, 60)
        pdf.multi_cell(140, 5.5, rec)
    
    # Footer note
    pdf.set_y(175)
    pdf.set_fill_color(248, 250, 252)
    pdf.rect(15, 175, 180, 20, 'F')
    pdf.set_xy(20, 180)
    pdf.set_font('Helvetica', 'I', 9)
    pdf.set_text_color(100, 110, 120)
    pdf.multi_cell(170, 5, "These recommendations are generated by AI based on your data patterns. Review and adapt them to your specific business context.")
    
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    base_name = Path(input_filename).stem
    path = f"{OUTPUT_DIR}/{base_name}_{ts}.pdf"
    pdf.output(path)
    return path


def generate_pptx(data, input_filename):
    """Build a stunning, executive-ready PowerPoint presentation"""
    from pptx.enum.shapes import MSO_SHAPE
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    report_title = Path(input_filename).stem.replace('_', ' ').replace('-', ' ').title()
    kpis = data['kpis']
    
    # Calculate top performer for insights
    top_campaign = data['campaigns'].row(0, named=True)
    top_roi = ((top_campaign['revenue'] - top_campaign['spend']) / top_campaign['spend'] * 100) if top_campaign['spend'] > 0 else 0
    
    # Color palette
    DARK_BLUE = RGBColor(20, 40, 80)
    ACCENT_BLUE = RGBColor(46, 134, 171)
    LIGHT_BLUE = RGBColor(230, 244, 250)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(50, 55, 65)
    LIGHT_GRAY = RGBColor(245, 247, 250)
    SUCCESS_GREEN = RGBColor(40, 167, 69)
    ALERT_RED = RGBColor(220, 53, 69)
    
    # ========== SLIDE 1: Title Slide ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Full dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_BLUE
    accent.line.fill.background()
    
    # Left accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.08), Inches(2))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_BLUE
    stripe.line.fill.background()
    
    # Main title
    title = slide.shapes.add_textbox(Inches(1.1), Inches(1.5), Inches(10), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = report_title
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(10), Inches(0.6))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Analytics Report"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(180, 200, 220)
    
    # Date range
    date_box = slide.shapes.add_textbox(Inches(1.1), Inches(3.3), Inches(10), Inches(0.4))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{data['dates'][0]}  |  {data['dates'][-1]}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(140, 160, 180)
    
    # Hero KPI cards at bottom
    kpi_highlights = [
        ('TOTAL REVENUE', f"${kpis['revenue']:,.0f}", SUCCESS_GREEN),
        ('RETURN ON INVESTMENT', f"{kpis['roi']:.1f}%", ACCENT_BLUE if kpis['roi'] > 0 else ALERT_RED),
        ('TOTAL IMPRESSIONS', f"{kpis['impressions']:,}", ACCENT_BLUE),
    ]
    
    for i, (label, value, accent_color) in enumerate(kpi_highlights):
        x = Inches(1 + i * 4)
        y = Inches(4.8)
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 50, 90)
        card.line.fill.background()
        
        # Top accent line
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.6), Inches(0.06))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = accent_color
        accent_line.line.fill.background()
        
        # Label
        lbl = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(3.2), Inches(0.4))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(150, 170, 190)
        
        # Value
        val = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), Inches(3.2), Inches(0.8))
        tf = val.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(11), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}  |  Powered by TrendSpotter AI"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 120, 140)
    
    # ========== SLIDE 2: KPI Overview ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Performance Indicators"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # KPI Grid - 2 rows x 3 columns
    kpi_data = [
        ('IMPRESSIONS', f"{kpis['impressions']:,}", 'Total ad views delivered', ACCENT_BLUE),
        ('CLICKS', f"{kpis['clicks']:,}", f"Click-through rate: {kpis['ctr']:.2f}%", RGBColor(162, 59, 114)),
        ('CONVERSIONS', f"{kpis['conversions']:,}", 'Completed actions', SUCCESS_GREEN),
        ('REVENUE', f"${kpis['revenue']:,.0f}", 'Total earnings generated', SUCCESS_GREEN),
        ('AD SPEND', f"${kpis['spend']:,.0f}", 'Marketing investment', RGBColor(255, 193, 7)),
        ('ROI', f"{kpis['roi']:.1f}%", 'Return on investment', SUCCESS_GREEN if kpis['roi'] > 50 else ALERT_RED),
    ]
    
    for i, (label, value, subtitle, color) in enumerate(kpi_data):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.6 + row * 2.6)
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()
        
        # Color accent on left
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(2.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        
        # Label
        lbl = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.4))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 105, 115)
        
        # Value
        val = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.8), Inches(3.3), Inches(0.9))
        tf = val.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        
        # Subtitle
        sub = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.7), Inches(3.3), Inches(0.4))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(120, 125, 135)
    
    # Key insight box at bottom
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = LIGHT_BLUE
    insight_box.line.fill.background()
    
    insight = slide.shapes.add_textbox(Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.5))
    tf = insight.text_frame
    p = tf.paragraphs[0]
    p.text = f"KEY INSIGHT: {top_campaign['campaign_name']} leads with ${top_campaign['revenue']:,.0f} revenue and {top_roi:.0f}% ROI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # ========== SLIDE 3: Executive Summary ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # AI badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(2.2), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = LIGHT_BLUE
    badge.line.fill.background()
    
    badge_text = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(2), Inches(0.35))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Generated Analysis"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Summary content box
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.8))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = LIGHT_GRAY
    summary_box.line.fill.background()
    
    summary = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.4))
    tf = summary.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data['summary']
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.line_spacing = 1.3
    
    # Anomalies section
    if data['anomalies']:
        anom_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.5))
        anom_header.fill.solid()
        anom_header.fill.fore_color.rgb = RGBColor(255, 235, 235)
        anom_header.line.fill.background()
        
        anom_title = slide.shapes.add_textbox(Inches(0.9), Inches(5.3), Inches(10), Inches(0.4))
        tf = anom_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"ANOMALIES DETECTED: {len(data['anomalies'])} unusual patterns identified"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ALERT_RED
        
        for i, a in enumerate(data['anomalies'][:3]):
            anom = slide.shapes.add_textbox(Inches(0.9), Inches(5.85 + i * 0.45), Inches(11), Inches(0.4))
            tf = anom.text_frame
            p = tf.paragraphs[0]
            p.text = f"  -  {a['desc'][:75]}"
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(100, 70, 70)
    
    # ========== SLIDE 4: Performance Dashboard ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Dashboard"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if os.path.exists(data['charts']['dashboard']):
        slide.shapes.add_picture(data['charts']['dashboard'], Inches(0.4), Inches(1.4), width=Inches(12.5))
    
    # ========== SLIDE 5: Campaign Performance ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Campaign Performance"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if os.path.exists(data['charts']['campaign']):
        slide.shapes.add_picture(data['charts']['campaign'], Inches(0.4), Inches(1.4), width=Inches(12.5))
    
    # ========== SLIDE 6: Strategic Recommendations ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Strategic Recommendations"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Recommendation cards with numbers
    rec_colors = [SUCCESS_GREEN, ACCENT_BLUE, RGBColor(162, 59, 114), RGBColor(255, 193, 7)]
    valid_recs = [r for r in data['recommendations'] if r and len(r.strip()) > 3]
    
    for i, rec in enumerate(valid_recs[:4]):
        y = Inches(1.5 + i * 1.4)
        color = rec_colors[i % 4]
        
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(12.1), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.fill.background()
        
        # Left accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.1), Inches(1.15))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()
        
        # Number circle
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y + Inches(0.25), Inches(0.65), Inches(0.65))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        num_circle.line.fill.background()
        
        num = slide.shapes.add_textbox(Inches(1), y + Inches(0.32), Inches(0.65), Inches(0.5))
        tf = num.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Recommendation text
        txt = slide.shapes.add_textbox(Inches(1.9), y + Inches(0.3), Inches(10.5), Inches(0.7))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rec
        p.font.size = Pt(17)
        p.font.color.rgb = DARK_GRAY
    
    # Footer note
    note = slide.shapes.add_textbox(Inches(0.6), Inches(7), Inches(12), Inches(0.4))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Recommendations generated by AI based on data patterns. Review and adapt to your business context."
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = RGBColor(130, 135, 145)
    
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    base_name = Path(input_filename).stem
    path = f"{OUTPUT_DIR}/{base_name}_{ts}.pptx"
    prs.save(path)
    return path


def main():
    start_time = time.time()
    
    output_format = 'pdf'
    input_file = INPUT_FILE
    
    args = sys.argv[1:]
    for arg in args:
        if arg in ['--pptx', '-p', 'pptx', 'ppt']:
            output_format = 'pptx'
        elif arg in ['--pdf', 'pdf']:
            output_format = 'pdf'
        elif arg in ['--help', '-h']:
            print("TrendSpotter - Automated Report Generator")
            print("\nUsage: python trendspotter.py [INPUT_FILE] [--pdf | --pptx]")
            print("\nSupports CSV, JSON, NDJSON, Parquet, and SQLite files.")
            print("Column names are auto-detected - just drop your export and go.")
            print("\nOptions:")
            print("  INPUT_FILE  Path to data file (default: data/sample/ad_performance.csv)")
            print("  --pdf       Generate PDF report (default)")
            print("  --pptx      Generate PowerPoint report")
            print("\nExamples:")
            print("  python trendspotter.py mydata.csv --pdf")
            print("  python trendspotter.py export.json --pptx")
            print("  python trendspotter.py database.sqlite --pdf")
            return
        elif arg.endswith(('.csv', '.json', '.ndjson', '.parquet', '.sqlite', '.db', '.sqlite3')):
            input_file = arg
    
    print("=" * 60)
    print(f"TRENDSPOTTER - Generating {output_format.upper()}")
    print("=" * 60)
    
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(CHARTS_DIR, exist_ok=True)
    
    # Load the data
    print(f"\n[1/6] Loading Data from {input_file}...")
    try:
        df, metadata = load_data(input_file)
    except FileNotFoundError as e:
        print(f"\nERROR: {e}")
        print("   Check the file path and try again.")
        sys.exit(1)
    except ValueError as e:
        print(f"\nERROR: {e}")
        sys.exit(1)
    except RuntimeError as e:
        print(f"\nERROR: {e}")
        sys.exit(1)
    
    # Show what we found in the data
    print(f"      Loaded {len(df)} records")
    print(f"\n      Data Understanding:")
    print(f"      - Original columns: {metadata.get('original_cols', [])}")
    if metadata.get('mapping'):
        print(f"      - Mapped to standard fields:")
        for orig, mapped in metadata['mapping'].items():
            print(f"           {orig} -> {mapped}")
    if metadata.get('synthetic'):
        print(f"      - Auto-generated fields: {metadata['synthetic']}")
    else:
        print(f"      - All required fields found in data")
    
    # Add derived metrics
    print("\n[2/6] Calculating Metrics...")
    if df['impressions'].sum() > 0:
        df = df.with_columns(
            (pl.col('clicks') / pl.col('impressions') * 100).fill_nan(0).round(2).alias('ctr')
        )
    else:
        df = df.with_columns(pl.lit(0.0).alias('ctr'))
    
    if df['spend'].sum() > 0:
        df = df.with_columns(
            ((pl.col('revenue') - pl.col('spend')) / pl.col('spend') * 100).fill_nan(0).round(2).alias('roi')
        )
    else:
        df = df.with_columns(pl.lit(0.0).alias('roi'))
    print("      CTR and ROI calculated")
    
    # Look for outliers in the data
    print("\n[3/6] Detecting Anomalies...")
    features = df.select(['impressions', 'clicks', 'spend', 'revenue']).to_numpy()
    anomalies = []
    if len(df) >= 10 and np.std(features) > 0:
        iso_forest = IsolationForest(contamination=0.1, random_state=42, n_estimators=50)
        predictions = iso_forest.fit_predict(features)
        df = df.with_columns(pl.Series('is_anomaly', predictions == -1))
        
        for row in df.filter(pl.col('is_anomaly')).iter_rows(named=True):
            anomalies.append({
                'campaign': row['campaign_name'],
                'date': str(row['date']),
                'region': row.get('region', 'Unknown'),
                'desc': f"{row['campaign_name']} - unusual metrics on {row['date']}"
            })
    else:
        df = df.with_columns(pl.lit(False).alias('is_anomaly'))
    print(f"      {len(anomalies)} anomalies found")
    
    # Roll up the numbers
    print("\n[4/6] Aggregating Data...")
    kpis = {
        'impressions': int(df['impressions'].sum()),
        'clicks': int(df['clicks'].sum()),
        'conversions': int(df['conversions'].sum()),
        'revenue': float(df['revenue'].sum()),
        'spend': float(df['spend'].sum()),
    }
    kpis['roi'] = ((kpis['revenue'] - kpis['spend']) / kpis['spend'] * 100) if kpis['spend'] > 0 else 0
    kpis['ctr'] = (kpis['clicks'] / kpis['impressions'] * 100) if kpis['impressions'] > 0 else 0
    
    campaigns = df.group_by('campaign_name').agg([
        pl.col('impressions').sum(),
        pl.col('clicks').sum(),
        pl.col('revenue').sum(),
        pl.col('spend').sum()
    ]).sort('revenue', descending=True)
    
    dates = sorted([str(d) for d in df['date'].unique().to_list()])
    print(f"      {len(campaigns)} campaigns, date range: {dates[0]} to {dates[-1]}")
    
    # Build the visualizations
    print("\n[5/6] Generating Charts...")
    charts = generate_charts(df, dates, campaigns, metadata)
    print("      Charts generated successfully")
    
    # Get AI insights
    print("\n[6/6] AI Analysis...")
    
    if API_KEY:
        client = OpenAI(api_key=API_KEY)
        
        context = f"""Data Analysis ({dates[0]} to {dates[-1]}):
- Impressions: {kpis['impressions']:,}, Clicks: {kpis['clicks']:,}, Conversions: {kpis['conversions']:,}
- Revenue: ${kpis['revenue']:,.0f}, Spend: ${kpis['spend']:,.0f}, ROI: {kpis['roi']:.1f}%
- Anomalies: {len(anomalies)}
Campaigns: {', '.join([f"{r['campaign_name']} (${r['revenue']:,.0f})" for r in campaigns.head(10).iter_rows(named=True)])}"""

        try:
            resp = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Senior Data Analyst. Write 100-word executive summary with key insights."},
                    {"role": "user", "content": context}
                ],
                max_tokens=200
            )
            summary = resp.choices[0].message.content
            
            resp2 = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Give exactly 4 actionable recommendations. Each must be under 80 characters. Format: one per line, no numbering."},
                    {"role": "user", "content": context}
                ],
                max_tokens=300
            )
            raw_recs = [l.strip().lstrip('0123456789.-) *') for l in resp2.choices[0].message.content.split('\n') if l.strip() and len(l) > 5]
            # Ensure recommendations are complete and not truncated
            recs = []
            for r in raw_recs[:4]:
                if len(r) > 80:
                    # Truncate at last space before 80 chars
                    r = r[:77].rsplit(' ', 1)[0] + '...'
                recs.append(r)
            # Pad with defaults if needed
            defaults = ["Optimize high-spend low-ROI campaigns", "Scale top performers", "Review anomalies", "A/B test creatives"]
            while len(recs) < 4:
                recs.append(defaults[len(recs)])
        except Exception as e:
            print(f"      Warning: AI analysis failed ({e})")
            summary = f"Data contains {len(df)} records across {len(campaigns)} campaigns with total revenue of ${kpis['revenue']:,.0f}."
            recs = ["Review campaign performance", "Investigate anomalies", "Optimize underperforming campaigns", "Scale successful campaigns"]
    else:
        print("      No API key found - using basic summary")
        summary = f"Data contains {len(df)} records across {len(campaigns)} campaigns. Total revenue: ${kpis['revenue']:,.0f}, Total spend: ${kpis['spend']:,.0f}, ROI: {kpis['roi']:.1f}%."
        recs = ["Review campaign performance", "Investigate anomalies", "Optimize underperforming campaigns", "Scale successful campaigns"]
    
    print("      Analysis complete")
    
    report_data = {
        'kpis': kpis,
        'dates': dates,
        'campaigns': campaigns,
        'anomalies': anomalies,
        'summary': summary,
        'recommendations': recs,
        'charts': charts
    }
    
    # Build the final report
    print(f"\nGenerating {output_format.upper()} report...")
    try:
        if output_format == 'pptx':
            path = generate_pptx(report_data, input_file)
        else:
            path = generate_pdf(report_data, input_file)
    except Exception as e:
        print(f"\nERROR generating report: {e}")
        sys.exit(1)
    
    elapsed = time.time() - start_time
    print(f"\n{'='*60}")
    print(f"DONE in {elapsed:.1f} seconds!")
    print(f"Output: {path}")
    print(f"{'='*60}")
    
    return path


if __name__ == "__main__":
    main()
